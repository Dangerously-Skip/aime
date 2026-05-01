#!/usr/bin/env python3
"""
Generic Markdown to PowerPoint Converter

Converts markdown documents to PowerPoint presentations using a template
and configuration file. Supports a domain-specific language (DSL) for
controlling slide layouts and styling.

Usage:
    python3 markdown_to_pptx.py input.md --template Template.pptx --output output.pptx
    python3 markdown_to_pptx.py input.md --config pptx_config.yaml

See MARKDOWN_TO_PPTX_DSL.md for full documentation.
"""

import re
import yaml
import argparse
import tempfile
import subprocess
import io
import json
import xml.etree.ElementTree as ET
from pathlib import Path
from urllib.parse import urlparse, unquote
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

class SlideMetadata:
    """Metadata for a single slide."""
    def __init__(self):
        self.slide_type = 'content'
        self.layout = None
        self.notes = ''
        self.transition = None
        self.subtitle = ''
        self.authors = ''
        self.date = ''
        self.image = None
        self.html_source = None  # For editable architecture diagrams

class MarkdownSlide:
    """Represents a single slide parsed from markdown."""
    def __init__(self, title, content, metadata):
        self.title = title
        self.content = content
        self.metadata = metadata

class MarkdownToPptxConverter:
    """Converts markdown to PowerPoint using template and config."""

    def __init__(self, template_path, config_path='pptx_config.yaml', brand_config_path='brand_config.yaml'):
        """Initialize converter with template and configuration."""
        self.template_path = Path(template_path)
        self.config_path = Path(config_path)
        self.brand_config_path = Path(brand_config_path)
        self.markdown_dir = None  # Will be set during convert()

        # Load configuration
        with open(self.config_path, 'r') as f:
            self.config = yaml.safe_load(f)

        # Load brand configuration
        if not self.brand_config_path.is_absolute():
            self.brand_config_path = self.config_path.parent / self.brand_config_path

        if self.brand_config_path.exists():
            with open(self.brand_config_path, 'r') as f:
                self.brand_config = yaml.safe_load(f)
        else:
            print(f"Warning: Brand config not found at {self.brand_config_path}, using defaults")
            self.brand_config = None

        # Override template path if specified in config
        if self.config.get('template', {}).get('path'):
            config_template_path = Path(self.config['template']['path'])
            # Resolve relative to config file location if not absolute
            if not config_template_path.is_absolute():
                config_template_path = self.config_path.parent / config_template_path
            self.template_path = config_template_path

        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")

    def parse_metadata_comment(self, line):
        """Parse metadata from HTML comment."""
        match = re.match(r'<!--\s*(\w+):\s*(.+?)\s*-->', line)
        if match:
            key = match.group(1)
            value = match.group(2)
            return key, value
        return None, None

    def parse_slide(self, slide_text):
        """Parse a single slide from markdown text."""
        lines = slide_text.strip().split('\n')
        if not lines:
            return None

        # Parse slide header
        header_line = lines[0]
        match = re.match(r'##\s*SLIDE\s*(?:(\w+):)?\s*(.+)', header_line)
        if not match:
            return None

        slide_type = match.group(1) or 'content'
        title = match.group(2).strip()

        # Parse metadata comments
        metadata = SlideMetadata()
        metadata.slide_type = slide_type

        content_start = 1
        for i, line in enumerate(lines[1:], 1):
            if line.strip().startswith('<!--'):
                key, value = self.parse_metadata_comment(line)
                if key:
                    setattr(metadata, key, value)
                content_start = i + 1
            else:
                break

        # Get content
        content = lines[content_start:]

        return MarkdownSlide(title, content, metadata)

    def parse_markdown(self, md_path):
        """Parse markdown file into slides."""
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Split by slide markers
        slide_pattern = r'##\s*SLIDE\s*(?:\w+:)?\s*.+?(?=##\s*SLIDE|$)'
        slide_texts = re.findall(slide_pattern, content, re.DOTALL)

        slides = []
        for slide_text in slide_texts:
            slide = self.parse_slide(slide_text)
            if slide:
                slides.append(slide)

        return slides

    def clean_markdown_text(self, text):
        """Remove markdown formatting from text."""
        if self.config['cleaning']['remove_markdown_formatting']:
            # Remove bold/italic
            text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\*(.+?)\*', r'\1', text)
            text = re.sub(r'__(.+?)__', r'\1', text)
            text = re.sub(r'_(.+?)_', r'\1', text)

        if self.config['cleaning']['remove_links']:
            # Remove markdown links but keep text
            text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)

        if self.config['cleaning']['remove_emojis']:
            # Remove common emojis
            text = re.sub(r'[🏆💰📈🚀✅📊💡🎯🔄💪🌟⚡🎁🤖🔌🧠💂📖🎖️🧑‍🏫👨‍🎓🚢🗺️👥🔥💵📱🌐🔒⏱️✨]+\s*', '', text)

        # Remove headers
        text = re.sub(r'^#+\s+', '', text)

        return text.strip()

    def apply_text_with_hyperlinks(self, paragraph, text, slide_part):
        """Apply text to paragraph with hyperlink support.

        Parses markdown links [text](url) and bold **text** formatting.
        """
        from pptx.oxml import parse_xml

        # Clear existing paragraph content
        paragraph.clear()

        # Process text with both links and bold formatting
        # First, find all markdown links and bold segments
        link_pattern = r'\[([^\]]+)\]\(([^\)]+)\)'
        bold_pattern = r'\*\*(.+?)\*\*'

        # Combine patterns to find all special formatting
        combined_pattern = f'({link_pattern}|{bold_pattern})'

        last_end = 0
        has_formatting = False

        for match in re.finditer(combined_pattern, text):
            has_formatting = True

            # Add text before the formatted segment
            before_text = text[last_end:match.start()]
            if before_text:
                run = paragraph.add_run()
                run.text = before_text

            # Check if this is a link or bold
            if match.group(2) and match.group(3):  # Link: [text](url)
                link_text = match.group(2)
                link_url = match.group(3)

                # Ensure URL has protocol
                if not link_url.startswith(('http://', 'https://')):
                    link_url = 'https://' + link_url

                # Create run with hyperlink
                run = paragraph.add_run()
                run.text = link_text

                try:
                    # Create external relationship
                    rel_id = slide_part.relate_to(
                        link_url,
                        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
                        is_external=True
                    )

                    # Add hyperlink XML
                    rPr = run._r.get_or_add_rPr()
                    hlink_xml = f'''<a:hlinkClick xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                                                  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                                                  r:id="{rel_id}"/>'''
                    hlink = parse_xml(hlink_xml)
                    rPr.append(hlink)

                    # Style as hyperlink
                    run.font.color.rgb = RGBColor(5, 99, 193)
                    run.font.underline = True
                except Exception as e:
                    print(f"Warning: Failed to create hyperlink for '{link_text}': {e}")

            elif match.group(4):  # Bold: **text**
                bold_text = match.group(4)
                run = paragraph.add_run()
                run.text = bold_text
                run.font.bold = True

            last_end = match.end()

        # Add remaining text after last formatted segment
        remaining_text = text[last_end:]
        if remaining_text:
            run = paragraph.add_run()
            run.text = remaining_text

        # If no formatting found, just set text normally
        if not has_formatting:
            paragraph.text = text

    def auto_detect_slide_type(self, slide):
        """Auto-detect slide type based on content."""
        content = slide.content
        title = slide.title.lower()

        # Check for table
        pipe_lines = [l for l in content if '|' in l]
        if len(pipe_lines) >= self.config['auto_detect']['table']['min_pipe_lines']:
            return 'table'

        # Check for metrics
        for pattern in self.config['auto_detect']['metrics']['patterns']:
            if any(re.search(pattern, line) for line in content[:5]):
                return 'metrics'

        # Check for quote
        for pattern in self.config['auto_detect']['quote']['patterns']:
            if any(re.search(pattern, line) for line in content[:3]):
                return 'quote'

        # Check for section header
        for keyword in self.config['auto_detect']['section']['keywords']:
            if keyword in title:
                return 'section'

        # Check for two column (if very long)
        if len(content) > self.config['auto_detect']['two_column']['min_lines']:
            return 'two_column'

        return 'content'

    def find_layout_by_name(self, prs, layout_name):
        """Find a layout by name across all slide masters.

        Args:
            prs: Presentation object
            layout_name: Name of layout to find (matches layout names in the template)

        Returns:
            SlideLayout object if found, None otherwise
        """
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name == layout_name:
                    return layout
        return None

    def get_layout_index(self, slide):
        """Get the layout index for a slide."""
        # Use explicit layout if specified
        if slide.metadata.layout:
            layout_name = slide.metadata.layout
            if layout_name in self.config['layouts']:
                return self.config['layouts'][layout_name]

        # Auto-detect if needed
        if slide.metadata.slide_type == 'content':
            detected_type = self.auto_detect_slide_type(slide)
            slide.metadata.slide_type = detected_type

        # Get default layout for slide type
        slide_type = slide.metadata.slide_type
        if slide_type in self.config['slide_type_defaults']:
            layout_name = self.config['slide_type_defaults'][slide_type]
            if layout_name in self.config['layouts']:
                return self.config['layouts'][layout_name]

        # Fallback to title_and_content
        return self.config['layouts']['title_and_content']

    def add_icon_to_slide(self, pptx_slide, slide):
        """Add icon(s) to a slide if specified in metadata."""
        # Check for single icon
        if hasattr(slide.metadata, 'icon') and slide.metadata.icon:
            icon_path = Path(slide.metadata.icon)

            # Try alternate extensions if file not found
            if not icon_path.exists():
                # Try .gif if .png specified, or vice versa
                if icon_path.suffix == '.png':
                    icon_path = icon_path.with_suffix('.gif')
                elif icon_path.suffix == '.gif':
                    icon_path = icon_path.with_suffix('.png')

            if not icon_path.exists():
                print(f"Warning: Icon not found: {slide.metadata.icon}")
                return

            # Position icon in top-right corner (small size)
            left = Inches(8.5)  # Right side
            top = Inches(0.5)   # Top
            height = Inches(0.6)  # Small icon

            try:
                pptx_slide.shapes.add_picture(str(icon_path), left, top, height=height)
            except Exception as e:
                print(f"Warning: Could not add icon {icon_path}: {e}")

        # Check for multiple icons
        if hasattr(slide.metadata, 'icons') and slide.metadata.icons:
            icon_paths = [p.strip() for p in slide.metadata.icons.split(',')]

            # Position multiple icons horizontally near title
            start_left = Inches(7.0)
            top = Inches(0.5)
            height = Inches(0.5)
            spacing = Inches(0.6)

            for i, icon_str in enumerate(icon_paths):
                icon_path = Path(icon_str.strip())

                # Try alternate extensions if file not found
                if not icon_path.exists():
                    if icon_path.suffix == '.png':
                        icon_path = icon_path.with_suffix('.gif')
                    elif icon_path.suffix == '.gif':
                        icon_path = icon_path.with_suffix('.png')

                if not icon_path.exists():
                    print(f"Warning: Icon not found: {icon_str}")
                    continue

                left = start_left + (i * spacing)
                try:
                    pptx_slide.shapes.add_picture(str(icon_path), left, top, height=height)
                except Exception as e:
                    print(f"Warning: Could not add icon {icon_path}: {e}")

    def create_title_slide(self, prs, slide):
        """Create a title slide.

        Searches for a branded title layout named "Title True Green" in the template.
        Falls back to the configured layout index if not found.
        """
        # Search for branded title layout in all slide masters
        title_layout = None
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name == "Title True Green":
                    title_layout = layout
                    break
            if title_layout:
                break

        # Use found layout or fallback
        if title_layout:
            pptx_slide = prs.slides.add_slide(title_layout)
        else:
            # Fallback to configured layout
            layout_idx = self.get_layout_index(slide)
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Fill placeholders if they exist
        title_filled = False
        subtitle_filled = False

        for shape in pptx_slide.shapes:
            if not shape.is_placeholder:
                continue

            # Title placeholder
            if shape.placeholder_format.type == 1:  # TITLE
                shape.text = self.clean_markdown_text(slide.title)
                # Bold the title
                if shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].font.bold = True
                title_filled = True

            # Body/Subtitle placeholder
            elif shape.placeholder_format.type == 2:  # BODY
                # Build subtitle text with all metadata
                subtitle_parts = []
                if slide.metadata.subtitle:
                    subtitle_parts.append(slide.metadata.subtitle)
                if slide.metadata.authors:
                    subtitle_parts.append(slide.metadata.authors)
                if slide.metadata.date:
                    subtitle_parts.append(slide.metadata.date)

                if subtitle_parts:
                    shape.text = '\n'.join(subtitle_parts)
                subtitle_filled = True

        # If no placeholders found, add text boxes manually (fallback)
        if not title_filled:
            # Manual text placement for template with no placeholders
            # Position text on left side (dark green area), away from right logo
            left_margin = Inches(0.8)
            content_width = Inches(7.5)  # Keep text in left 60% of slide

            # Title (large, bold, white)
            title_top = Inches(2.2)
            title_box = pptx_slide.shapes.add_textbox(left_margin, title_top, content_width, Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.text = self.clean_markdown_text(slide.title)
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)  # White

            # Subtitle (medium, white)
            if slide.metadata.subtitle:
                subtitle_top = Inches(3.5)
                subtitle_box = pptx_slide.shapes.add_textbox(left_margin, subtitle_top, content_width, Inches(0.7))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide.metadata.subtitle
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = RGBColor(255, 255, 255)  # White

            # Authors (smaller, white)
            if slide.metadata.authors:
                authors_top = Inches(4.4)
                authors_box = pptx_slide.shapes.add_textbox(left_margin, authors_top, content_width, Inches(0.5))
                authors_frame = authors_box.text_frame
                authors_frame.text = slide.metadata.authors
                authors_para = authors_frame.paragraphs[0]
                authors_para.font.size = Pt(18)
                authors_para.font.color.rgb = RGBColor(255, 255, 255)  # White

            # Date (smaller, white)
            if slide.metadata.date:
                date_top = Inches(5.0)
                date_box = pptx_slide.shapes.add_textbox(left_margin, date_top, content_width, Inches(0.4))
                date_frame = date_box.text_frame
                date_frame.text = slide.metadata.date
                date_para = date_frame.paragraphs[0]
                date_para.font.size = Pt(16)
                date_para.font.color.rgb = RGBColor(255, 255, 255)  # White
        else:
            # Old template behavior with placeholders
            for shape in pptx_slide.placeholders:
                if shape.placeholder_format.type == 7:  # OBJECT type
                    text_frame = shape.text_frame
                    text_frame.clear()

                    # Add subtitle
                    if slide.metadata.subtitle:
                        p = text_frame.paragraphs[0]
                        p.text = slide.metadata.subtitle
                        p.font.size = Pt(24)

                    # Add authors
                    if slide.metadata.authors:
                        p = text_frame.add_paragraph()
                        p.text = slide.metadata.authors
                        p.font.size = Pt(18)
                        p.space_before = Pt(20)

                    # Add date
                    if slide.metadata.date:
                        p = text_frame.add_paragraph()
                        p.text = slide.metadata.date
                        p.font.size = Pt(16)
                        p.space_before = Pt(10)

                    break

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def parse_markdown_table(self, content_lines):
        """Parse markdown table into rows and columns."""
        table_lines = [l for l in content_lines if '|' in l]
        if len(table_lines) < 2:
            return None

        rows = []
        for line in table_lines:
            # Skip separator line (e.g., |---|---|)
            if re.match(r'^\s*\|[\s\-:|]+\|\s*$', line):
                continue

            # Split by | and clean up
            cells = [cell.strip() for cell in line.split('|')]
            # Remove empty first/last if line starts/ends with |
            cells = [c for c in cells if c]

            if cells:
                rows.append(cells)

        return rows if len(rows) >= 2 else None

    def create_table_slide(self, prs, slide):
        """Create a slide with a table."""
        layout_idx = self.get_layout_index(slide)
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if pptx_slide.shapes.title:
            pptx_slide.shapes.title.text = self.clean_markdown_text(slide.title)

        # Parse table from content
        table_data = self.parse_markdown_table(slide.content)

        if not table_data:
            # Fallback to content slide if no valid table
            return self.create_content_slide(prs, slide)

        # Find content placeholder and create table
        for shape in pptx_slide.placeholders:
            if shape.placeholder_format.type == 7:  # OBJECT type
                # Remove the placeholder shape
                sp = shape.element
                sp.getparent().remove(sp)

                # Calculate table dimensions
                rows_count = len(table_data)
                cols_count = len(table_data[0]) if table_data else 0

                # Position table where placeholder was
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Create table
                table_shape = pptx_slide.shapes.add_table(
                    rows_count, cols_count, left, top, width, height
                )
                table = table_shape.table

                # Get brand colors if available
                if self.brand_config:
                    header_bg = self.brand_config['powerpoint_mapping']['rgb_colors']['brand_green']
                    header_fg = self.brand_config['powerpoint_mapping']['rgb_colors']['white']
                    header_size = self.brand_config['components']['tables']['header_font_size']
                    body_size = self.brand_config['components']['tables']['body_font_size']
                else:
                    # Fallback colors
                    header_bg = [20, 74, 56]
                    header_fg = [255, 255, 255]
                    header_size = 14
                    body_size = 12

                # Fill table with data
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_text in enumerate(row_data):
                        if col_idx < cols_count:
                            cell = table.cell(row_idx, col_idx)
                            # Clean markdown formatting from cell text
                            cell_text = self.clean_markdown_text(cell_text)
                            cell.text = cell_text

                            # Format header row with brand colors
                            if row_idx == 0:
                                cell.text_frame.paragraphs[0].font.bold = True
                                cell.text_frame.paragraphs[0].font.size = Pt(header_size)
                                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*header_fg)
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(*header_bg)
                            else:
                                cell.text_frame.paragraphs[0].font.size = Pt(body_size)

                break

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def create_content_slide(self, prs, slide):
        """Create a standard content slide.

        Can use custom layouts like "Title Warm White", "Title Sage Green", etc.
        Falls back to configured default layout.
        """
        # Try to use specified layout name first
        content_layout = None
        if slide.metadata.layout:
            content_layout = self.find_layout_by_name(prs, slide.metadata.layout)
            if not content_layout:
                print(f"Warning: Layout '{slide.metadata.layout}' not found for content slide")

        # Fall back to default content layout
        if not content_layout:
            layout_idx = self.get_layout_index(slide)
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        else:
            pptx_slide = prs.slides.add_slide(content_layout)

        # Set title
        if pptx_slide.shapes.title:
            pptx_slide.shapes.title.text = self.clean_markdown_text(slide.title)

        # Add content - use FIRST OBJECT placeholder only, delete any extras
        filled_placeholder = False
        placeholders_to_remove = []

        for shape in pptx_slide.placeholders:
            if shape.placeholder_format.type == 7:  # OBJECT type
                if not filled_placeholder:
                    # Fill the FIRST placeholder
                    text_frame = shape.text_frame
                    text_frame.clear()

                    max_bullets = self.config['limits']['max_bullets_per_slide']
                    content_lines = slide.content[:max_bullets]

                    first_content = True
                    for i, line in enumerate(content_lines):
                        line = line.strip()
                        if not line or line == '---':
                            continue

                        # Don't clean markdown text yet - preserve links for parsing
                        raw_line = line.strip()

                        if first_content:
                            p = text_frame.paragraphs[0]
                            first_content = False
                        else:
                            p = text_frame.add_paragraph()

                        # Detect indentation level and strip bullet markers
                        from pptx.util import Pt
                        is_bullet = False
                        if raw_line.startswith('- ') or raw_line.startswith('* '):
                            self.apply_text_with_hyperlinks(p, raw_line[2:], pptx_slide.part)
                            p.level = 0
                            is_bullet = True
                            # Reduce bullet font size
                            for run in p.runs:
                                run.font.size = Pt(14)
                        elif raw_line.startswith('  - ') or raw_line.startswith('  * '):
                            self.apply_text_with_hyperlinks(p, raw_line[4:], pptx_slide.part)
                            p.level = 1
                            is_bullet = True
                            for run in p.runs:
                                run.font.size = Pt(14)
                        elif raw_line.startswith('    - ') or raw_line.startswith('    * '):
                            self.apply_text_with_hyperlinks(p, raw_line[6:], pptx_slide.part)
                            p.level = 2
                            is_bullet = True
                            for run in p.runs:
                                run.font.size = Pt(14)
                        else:
                            # Non-bullet paragraph - reduce font size slightly and turn off bold
                            self.apply_text_with_hyperlinks(p, raw_line, pptx_slide.part)
                            for run in p.runs:
                                run.font.size = Pt(14)  # Smaller than default 18pt
                                run.font.bold = False  # Don't bold paragraph text

                        # Apply body font settings from config
                        if 'fonts' in self.config and 'body' in self.config['fonts']:
                            body_config = self.config['fonts']['body']
                            if 'bold' in body_config and is_bullet:
                                p.font.bold = body_config['bold']

                    filled_placeholder = True
                else:
                    # Mark EXTRA placeholders for removal
                    placeholders_to_remove.append(shape)

        # Remove extra placeholders
        for shape in placeholders_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def create_section_header_slide(self, prs, slide):
        """Create a section header slide.

        Uses template layouts like "Divider Bright Green", "Divider True Green", etc.
        Falls back to configured section_header layout if no specific layout requested.
        """
        # Try to use specified layout name first
        section_layout = None
        if slide.metadata.layout:
            section_layout = self.find_layout_by_name(prs, slide.metadata.layout)
            if not section_layout:
                print(f"Warning: Layout '{slide.metadata.layout}' not found for section slide")

        # Fall back to default section header layout
        if not section_layout:
            layout_idx = self.get_layout_index(slide)
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        else:
            pptx_slide = prs.slides.add_slide(section_layout)

        # Fill title placeholder if it exists
        if pptx_slide.shapes.title:
            pptx_slide.shapes.title.text = self.clean_markdown_text(slide.title)

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def create_two_column_slide(self, prs, slide):
        """Create a two-column slide."""
        layout_idx = self.get_layout_index(slide)
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if pptx_slide.shapes.title:
            pptx_slide.shapes.title.text = self.clean_markdown_text(slide.title)

        # Split content
        content = [l for l in slide.content if l.strip() and l.strip() != '---']

        # Check for manual column markers
        left_content = []
        right_content = []
        current_column = None

        for line in content:
            if line.strip() == '::: column-left':
                current_column = 'left'
            elif line.strip() == '::: column-right':
                current_column = 'right'
            elif line.strip() == ':::':
                current_column = None
            elif current_column == 'left':
                left_content.append(line)
            elif current_column == 'right':
                right_content.append(line)
            elif current_column is None:
                # No manual split, use auto-split
                if not left_content and not right_content:
                    # Sequential split
                    mid = len(content) // 2
                    left_content = content[:mid]
                    right_content = content[mid:]
                    break

        # Find content placeholders
        content_shapes = [s for s in pptx_slide.placeholders if s.placeholder_format.type == 7]

        if len(content_shapes) >= 2:
            # Left column
            self._fill_text_frame(content_shapes[0].text_frame, left_content)

            # Right column
            self._fill_text_frame(content_shapes[1].text_frame, right_content)

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def _fill_text_frame(self, text_frame, content_lines):
        """Fill a text frame with content."""
        text_frame.clear()

        for i, line in enumerate(content_lines[:10]):
            line = self.clean_markdown_text(line.strip())
            if not line:
                continue

            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # Handle bullets
            if line.startswith('- ') or line.startswith('* '):
                p.text = line[2:]
                p.level = 0
            else:
                p.text = line

    def generate_editable_architecture_slide(self, html_path):
        """Generate editable PowerPoint slide from HTML architecture diagram.

        Uses HTML → SVG → PPTX pipeline to create fully editable shapes.

        Args:
            html_path: Path to HTML architecture diagram file

        Returns:
            Path to generated PPTX file (in temp directory)
        """
        html_path = Path(html_path)
        if not html_path.exists():
            print(f"Warning: HTML source not found: {html_path}")
            return None

        # Create temp directory for intermediate files
        temp_dir = Path(tempfile.mkdtemp(prefix='pptx_arch_'))
        svg_path = temp_dir / 'architecture.svg'
        pptx_path = temp_dir / 'editable_slide.pptx'

        # Get path to helper scripts (in same directory as this script)
        script_dir = Path(__file__).parent
        html_to_svg_script = script_dir / 'html_architecture_to_svg.py'
        svg_to_pptx_script = script_dir / 'svg_to_pptx_custom.py'

        if not html_to_svg_script.exists() or not svg_to_pptx_script.exists():
            print(f"Warning: Editable architecture scripts not found in {script_dir}")
            return None

        try:
            # Step 1: HTML → SVG (with images extracted to JSON)
            print(f"    Converting HTML to SVG: {html_path.name}")
            result = subprocess.run(
                ['python3', str(html_to_svg_script), str(html_path), '--output', str(svg_path)],
                capture_output=True,
                text=True,
                timeout=30,
                cwd=html_path.parent
            )

            if result.returncode != 0 or not svg_path.exists():
                print(f"    Warning: HTML to SVG conversion failed")
                if result.stderr:
                    print(f"    Error: {result.stderr[:200]}")
                return None

            # Step 2: SVG → PPTX (with logos added)
            print(f"    Converting SVG to editable PPTX")
            result = subprocess.run(
                ['python3', str(svg_to_pptx_script), str(svg_path), '--output', str(pptx_path)],
                capture_output=True,
                text=True,
                timeout=60,
                cwd=html_path.parent
            )

            if result.returncode != 0 or not pptx_path.exists():
                print(f"    Warning: SVG to PPTX conversion failed")
                if result.stderr:
                    print(f"    Error: {result.stderr[:200]}")
                return None

            print(f"    ✓ Generated editable slide: {pptx_path}")
            return pptx_path

        except subprocess.TimeoutExpired:
            print(f"    Warning: Timeout generating editable slide from {html_path}")
            return None
        except Exception as e:
            print(f"    Warning: Error generating editable slide: {e}")
            return None

    def copy_shapes_from_slide(self, source_pptx_path, target_slide):
        """Copy all shapes from a source PPTX slide to a target slide.

        Args:
            source_pptx_path: Path to source PPTX file (should have 1 slide)
            target_slide: Target slide object to copy shapes into
        """
        try:
            source_prs = Presentation(str(source_pptx_path))

            if len(source_prs.slides) == 0:
                print("    Warning: Source PPTX has no slides")
                return False

            source_slide = source_prs.slides[0]

            # Copy all shapes using proper python-pptx methods
            shapes_copied = 0
            images_copied = 0
            autoshapes_copied = 0
            textboxes_copied = 0

            for shape in source_slide.shapes:
                # Skip title placeholders - we already have one
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                    continue

                try:
                    # Handle images first
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        # Get image stream and properties
                        image_part = shape.image.blob
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # Add image to target slide
                        target_slide.shapes.add_picture(
                            io.BytesIO(image_part),
                            left=left,
                            top=top,
                            width=width,
                            height=height
                        )
                        images_copied += 1
                        shapes_copied += 1

                    # Handle shapes (rectangles, etc.) BEFORE text boxes
                    # AutoShapes also have text_frame, so check them first
                    elif shape.shape_type == 1:  # AUTO_SHAPE (rectangles, etc.)
                        from pptx.enum.shapes import MSO_SHAPE
                        from pptx.util import Emu

                        # Get shape properties
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # Get the AutoShape type (rectangle = 1)
                        autoshape_type = shape.auto_shape_type

                        # Add shape to target slide
                        new_shape = target_slide.shapes.add_shape(
                            autoshape_type,
                            left, top, width, height
                        )

                        # Copy fill properties
                        if shape.fill.type == 1:  # Solid fill
                            new_shape.fill.solid()
                            if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                                new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                        elif shape.fill.type == 0:  # No fill
                            new_shape.fill.background()

                        # Copy line/border properties
                        if shape.line.width:
                            new_shape.line.width = shape.line.width
                        if hasattr(shape.line, 'color') and hasattr(shape.line.color, 'rgb'):
                            if shape.line.color.rgb:
                                new_shape.line.color.rgb = shape.line.color.rgb

                        autoshapes_copied += 1
                        shapes_copied += 1

                    # Handle text boxes (check this AFTER AutoShapes)
                    elif hasattr(shape, 'text_frame') and shape.has_text_frame:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # Create text box
                        new_textbox = target_slide.shapes.add_textbox(left, top, width, height)

                        # Copy text and formatting
                        new_textbox.text_frame.text = shape.text
                        if shape.text_frame.paragraphs:
                            for src_para, dst_para in zip(shape.text_frame.paragraphs, new_textbox.text_frame.paragraphs):
                                if src_para.runs:
                                    for src_run, dst_run in zip(src_para.runs, dst_para.runs):
                                        if src_run.font.size:
                                            dst_run.font.size = src_run.font.size
                                        if src_run.font.bold:
                                            dst_run.font.bold = src_run.font.bold
                                        if src_run.font.color.rgb:
                                            dst_run.font.color.rgb = src_run.font.color.rgb

                        textboxes_copied += 1
                        shapes_copied += 1

                    # Handle groups
                    elif shape.shape_type == 6:  # GROUP
                        # Groups are complex - use XML cloning as fallback
                        from copy import deepcopy
                        shape_elm = deepcopy(shape.element)
                        target_slide.shapes._spTree.append(shape_elm)
                        shapes_copied += 1

                except Exception as e:
                    print(f"    Warning: Could not copy shape type {shape.shape_type}: {str(e)[:100]}")
                    continue

            print(f"    ✓ Copied {shapes_copied} shapes: {autoshapes_copied} boxes, {textboxes_copied} text, {images_copied} images")
            return True

        except Exception as e:
            print(f"    Warning: Error copying shapes from {source_pptx_path}: {e}")
            return False

    def generate_shapes_from_svg_direct(self, svg_path, target_slide, presentation):
        """Generate shapes directly on target slide from SVG.

        This avoids coordinate system issues from copying between presentations.
        """
        svg_path = Path(svg_path)
        images_json_path = svg_path.with_suffix('.images.json')

        # Parse SVG
        tree = ET.parse(svg_path)
        root = tree.getroot()

        # Get SVG dimensions
        svg_width = float(root.get('width', 1333))
        svg_height = float(root.get('height', 750))

        # Get PowerPoint slide dimensions
        slide_width_inches = presentation.slide_width / 914400  # EMU to inches
        slide_height_inches = presentation.slide_height / 914400

        # Calculate scaling factors - limit diagram to 70% of slide width to leave room for text column
        max_diagram_width = slide_width_inches * 0.70
        max_diagram_height = slide_height_inches - 2.2  # Leave room for title/subtitle and bottom margin

        scale_x = max_diagram_width / svg_width
        scale_y = max_diagram_height / svg_height

        # Use uniform scaling (maintain aspect ratio)
        scale = min(scale_x, scale_y)
        scale_x = scale
        scale_y = scale

        # Offset for title area (push content down below title)
        title_offset_inches = 1.65

        print(f"    Generating shapes: SVG {svg_width}x{svg_height} → PPT {slide_width_inches:.2f}\"x{slide_height_inches:.2f}\"")

        # Helper to convert hex color
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            if len(hex_color) == 6:
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            return (0, 0, 0)

        # SVG namespace
        ns = {'svg': 'http://www.w3.org/2000/svg'}

        # Process rectangles
        shapes_added = 0
        for rect in root.findall('.//svg:rect', ns):
            try:
                x = float(rect.get('x', 0))
                y = float(rect.get('y', 0))
                width = float(rect.get('width', 0))
                height = float(rect.get('height', 0))

                if width < 2 or height < 2:
                    continue

                left = Inches(x * scale_x)
                top = Inches(y * scale_y + title_offset_inches)
                w = Inches(width * scale_x)
                h = Inches(height * scale_y)

                shape = target_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left, top, w, h
                )

                # Set fill
                fill = rect.get('fill', 'none')
                if fill and fill != 'none':
                    shape.fill.solid()
                    if fill.startswith('#'):
                        rgb = hex_to_rgb(fill)
                        shape.fill.fore_color.rgb = RGBColor(*rgb)
                else:
                    shape.fill.background()

                # Set border
                stroke = rect.get('stroke', 'none')
                stroke_width = rect.get('stroke-width', '0')

                if stroke and stroke != 'none':
                    shape.line.color.rgb = RGBColor(*hex_to_rgb(stroke))
                    try:
                        shape.line.width = Pt(float(stroke_width))
                    except:
                        shape.line.width = Pt(1)
                else:
                    shape.line.fill.background()

                shapes_added += 1
            except Exception as e:
                continue

        # Process text elements
        texts_added = 0
        for text_elem in root.findall('.//svg:text', ns):
            try:
                x = float(text_elem.get('x', 0))
                y = float(text_elem.get('y', 0))

                text_content = ''.join(text_elem.itertext()).strip()
                if not text_content:
                    continue

                svg_font_size = float(text_elem.get('font-size', 12))
                font_weight = text_elem.get('font-weight', 'normal')
                text_anchor = text_elem.get('text-anchor', 'start')

                pptx_font_size = max(6, svg_font_size * scale_y * 10)

                char_width = pptx_font_size * 0.7
                text_width_pts = len(text_content) * char_width * 1.5
                text_height_pts = pptx_font_size * 2.0

                text_width_inches = text_width_pts / 72
                text_height_inches = text_height_pts / 72

                left = Inches(x * scale_x)
                top = Inches(y * scale_y - text_height_inches / 2 + title_offset_inches)

                if text_anchor == 'middle':
                    left -= Inches(text_width_inches / 2)
                elif text_anchor == 'end':
                    left -= Inches(text_width_inches)

                w = Inches(text_width_inches)
                h = Inches(text_height_inches)

                textbox = target_slide.shapes.add_textbox(left, top, w, h)
                text_frame = textbox.text_frame
                text_frame.text = text_content
                text_frame.word_wrap = False

                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(pptx_font_size)
                    paragraph.font.bold = (font_weight in ['bold', '700', '600'])

                    fill = text_elem.get('fill', '#000000')
                    if fill.startswith('#'):
                        paragraph.font.color.rgb = RGBColor(*hex_to_rgb(fill))

                    if text_anchor == 'middle':
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif text_anchor == 'end':
                        paragraph.alignment = PP_ALIGN.RIGHT

                texts_added += 1
            except Exception as e:
                continue

        # Add images from JSON
        images_added = 0
        if images_json_path.exists():
            with open(images_json_path, 'r') as f:
                data = json.load(f)

            for img_data in data.get('images', []):
                try:
                    src = img_data['src']
                    x = img_data['x']
                    y = img_data['y']
                    width = img_data['width']
                    height = img_data['height']

                    if src.startswith('file://'):
                        img_path = unquote(urlparse(src).path)
                    else:
                        img_path = src

                    img_path = Path(img_path)
                    if not img_path.exists():
                        continue

                    left = Inches(x * scale_x)
                    top = Inches(y * scale_y + title_offset_inches)
                    w = Inches(width * scale_x)
                    h = Inches(height * scale_y)

                    target_slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
                    images_added += 1
                except Exception as e:
                    continue

        print(f"    ✓ Added {shapes_added} shapes, {texts_added} texts, {images_added} images")

        return True

    def create_image_slide(self, prs, slide):
        """Create a slide with a large image using the Picture with Caption layout."""
        # Use blank layout for more control
        blank_layout_idx = self.config['layouts']['blank']
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[blank_layout_idx])

        # Remove any default placeholders except title
        shapes_to_remove = []
        title_shape = None
        for shape in pptx_slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type == 1:  # TITLE
                    title_shape = shape
                else:
                    # Mark other placeholders for removal
                    shapes_to_remove.append(shape)

        # Remove non-title placeholders
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

        # Set title (hide if it's just a placeholder like ".")
        if title_shape:
            title_text = self.clean_markdown_text(slide.title)
            if title_text and title_text.strip() != '.':
                title_shape.text = title_text
            else:
                # Remove title placeholder for full-slide images
                sp = title_shape.element
                sp.getparent().remove(sp)
                title_shape = None

        # Add subtitle if specified
        subtitle_shape = None
        if hasattr(slide.metadata, 'subtitle') and slide.metadata.subtitle:
            # Add subtitle text box below title
            from pptx.util import Pt
            left = Inches(0.5)
            top = Inches(1.2)  # Below title
            width = prs.slide_width - Inches(1)
            height = Inches(0.4)

            subtitle_box = pptx_slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide.metadata.subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(16)
            subtitle_para.font.color.rgb = RGBColor(102, 102, 102)  # Gray
            subtitle_shape = subtitle_box

        # Get slide dimensions for positioning
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Check for editable architecture diagram (html_source)
        # This takes priority over static images
        html_source_path = None
        if hasattr(slide.metadata, 'html_source') and slide.metadata.html_source:
            html_source_path = Path(slide.metadata.html_source)
            # Make path relative to markdown file's directory if needed
            if not html_source_path.is_absolute():
                html_source_path = self.markdown_dir / html_source_path

            if html_source_path and html_source_path.exists():
                print(f"    🎨 Generating editable architecture diagram from HTML")
                editable_pptx_path = self.generate_editable_architecture_slide(html_source_path)

                if editable_pptx_path and editable_pptx_path.exists():
                    # Keep title placeholder - add title and subtitle
                    # (Commenting out title removal - we want the slide title to show)
                    # if title_shape:
                    #     sp = title_shape.element
                    #     sp.getparent().remove(sp)
                    #     title_shape = None

                    # Generate shapes directly on this slide instead of copying
                    # This avoids coordinate system issues from copying between presentations
                    svg_path = editable_pptx_path.parent / 'architecture.svg'
                    if svg_path.exists():
                        success = self.generate_shapes_from_svg_direct(svg_path, pptx_slide, prs)

                        if success:
                            # Check if there's additional slide content to add
                            if slide.content and len([line for line in slide.content if line.strip() and line.strip() != '---']) > 0:
                                # Add slide content as text on the right side
                                # Architecture diagram is ~70% width, text column is ~30% on right
                                from pptx.util import Pt
                                left = slide_width * 0.72  # Start at 72% to avoid overlap
                                top = Inches(1.8)  # Below subtitle
                                width = slide_width * 0.28 - Inches(0.3)  # 28% width - margins
                                height = slide_height - Inches(2.5)  # Leave room for title/subtitle and bottom

                                textbox = pptx_slide.shapes.add_textbox(left, top, width, height)
                                text_frame = textbox.text_frame
                                text_frame.word_wrap = True

                                for i, line in enumerate(slide.content):
                                    line = self.clean_markdown_text(line.strip())
                                    if not line or line == '---':
                                        continue

                                    if i == 0:
                                        p = text_frame.paragraphs[0]
                                    else:
                                        p = text_frame.add_paragraph()

                                    # Handle bold text for headers
                                    if line.startswith('**') and line.endswith('**'):
                                        p.text = line[2:-2]  # Remove ** markers
                                        p.font.bold = True
                                        p.font.size = Pt(16)
                                        p.font.color.rgb = RGBColor(255, 140, 0)  # Orange for DATA CONSUMERS
                                    elif line.startswith('- '):
                                        # Bullet point
                                        p.text = line[2:]  # Remove "- " prefix
                                        p.level = 1
                                        p.font.size = Pt(14)
                                    else:
                                        p.text = line
                                        p.font.size = Pt(14)

                            # Add icon if specified
                            self.add_icon_to_slide(pptx_slide, slide)
                            return pptx_slide
                        else:
                            print(f"    Warning: Failed to generate shapes, falling back to PNG")
                    else:
                        print(f"    Warning: SVG not found, falling back to PNG")
                else:
                    print(f"    Warning: Failed to generate editable slide, falling back to PNG")

        # Check for image in metadata
        image_path = None
        if hasattr(slide.metadata, 'image') and slide.metadata.image:
            image_path = Path(slide.metadata.image)
            # Make path relative to markdown file's directory if needed
            if not image_path.is_absolute():
                image_path = self.markdown_dir / image_path

        if image_path and image_path.exists():
            # Calculate image position and scale - full slide if no title, below title otherwise

            from PIL import Image as PILImage

            # Get image dimensions to calculate proper scaling
            try:
                with PILImage.open(image_path) as img:
                    img_width_px, img_height_px = img.size
                    img_aspect = img_width_px / img_height_px
            except:
                img_aspect = 16 / 9  # Default fallback

            if title_shape:
                # Position below title with more vertical space
                margin_left = Inches(0.5)
                margin_top = Inches(1.5)
                margin_right = Inches(0.5)
                margin_bottom = Inches(0.5)
            else:
                # Full slide - no title present, minimal margins
                margin_left = Inches(0.2)
                margin_top = Inches(0.2)
                margin_right = Inches(0.2)
                margin_bottom = Inches(0.2)

            available_width = slide_width - margin_left - margin_right
            available_height = slide_height - margin_top - margin_bottom

            # Calculate which dimension is the constraint
            slide_aspect = available_width / available_height

            if img_aspect > slide_aspect:
                # Image is wider - width is constraint
                width = available_width
                height = width / img_aspect
                left = margin_left
                top = margin_top + (available_height - height) / 2
            else:
                # Image is taller - height is constraint
                height = available_height
                width = height * img_aspect
                left = margin_left + (available_width - width) / 2
                top = margin_top

            try:
                pptx_slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
            except Exception as e:
                print(f"Warning: Could not add image {image_path}: {e}")
                # Add error message as text box
                from pptx.util import Pt
                textbox = pptx_slide.shapes.add_textbox(left, top, width, Inches(1))
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = f"Image not found: {image_path}"
                p.font.size = Pt(18)
        else:
            # No image found - add text content
            from pptx.util import Pt
            left = Inches(1)
            top = Inches(2)
            width = slide_width - Inches(2)
            height = slide_height - Inches(3)

            textbox = pptx_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for i, line in enumerate(slide.content[:5]):
                line = self.clean_markdown_text(line.strip())
                if not line or line == '---':
                    continue
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

        # Add icon if specified
        self.add_icon_to_slide(pptx_slide, slide)

        return pptx_slide

    def create_slide(self, prs, slide):
        """Create appropriate slide based on type."""
        slide_type = slide.metadata.slide_type

        if slide_type == 'title':
            return self.create_title_slide(prs, slide)
        elif slide_type == 'section':
            return self.create_section_header_slide(prs, slide)
        elif slide_type == 'two_column':
            return self.create_two_column_slide(prs, slide)
        elif slide_type == 'table':
            return self.create_table_slide(prs, slide)
        elif slide_type == 'image':
            return self.create_image_slide(prs, slide)
        else:
            # Default to content slide
            return self.create_content_slide(prs, slide)

    def convert(self, md_path, output_path):
        """Convert markdown file to PowerPoint presentation."""
        # Store markdown directory for resolving relative image paths
        self.markdown_dir = Path(md_path).resolve().parent

        print(f"\n{'='*80}")
        print(f"Converting: {md_path} → {output_path}")
        print(f"Template: {self.template_path}")
        print(f"{'='*80}\n")

        # Parse markdown
        slides = self.parse_markdown(md_path)
        print(f"Parsed {len(slides)} slides")

        # Load template
        prs = Presentation(self.template_path)

        # Remove any existing slides from template
        # Keep only the layouts and master slides, not content slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # Create slides
        for i, slide in enumerate(slides, 1):
            print(f"  [{i:2d}/{len(slides)}] {slide.title[:60]}... (type: {slide.metadata.slide_type})")
            self.create_slide(prs, slide)

        # Save
        prs.save(output_path)
        print(f"\n✓ Created: {output_path}")

def main():
    parser = argparse.ArgumentParser(
        description='Convert markdown to PowerPoint presentation',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.md --template Template.pptx --output output.pptx
  %(prog)s presentation.md --config pptx_config.yaml
        """
    )

    parser.add_argument('input', help='Input markdown file')
    parser.add_argument('--template', help='PowerPoint template file')
    parser.add_argument('--config', default='pptx_config.yaml',
                        help='Configuration file (default: pptx_config.yaml)')
    parser.add_argument('--output', help='Output PowerPoint file')

    args = parser.parse_args()

    # Determine template path
    template_path = args.template
    if not template_path:
        # Try to get from config
        try:
            with open(args.config, 'r') as f:
                config = yaml.safe_load(f)
                template_path = config.get('template', {}).get('path')
        except:
            pass

    if not template_path:
        parser.error('Template must be specified via --template or in config file')

    # Determine output path
    output_path = args.output
    if not output_path:
        input_path = Path(args.input)
        output_path = input_path.stem + '.pptx'

    # Convert
    converter = MarkdownToPptxConverter(template_path, args.config)
    converter.convert(args.input, output_path)

if __name__ == '__main__':
    main()
