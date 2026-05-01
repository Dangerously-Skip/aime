#!/usr/bin/env python3
"""
Custom SVG to PowerPoint converter with proper coordinate mapping.

Converts SVG shapes to native PowerPoint shapes with correct positioning.
Replaces svg2pptx which has coordinate system issues.
"""

import sys
import json
import argparse
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from urllib.parse import urlparse, unquote


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return (0, 0, 0)


def svg_to_pptx_custom(svg_path, output_path, slide_width_inches=13.33, slide_height_inches=7.5):
    """
    Convert SVG to PowerPoint with proper coordinate mapping.

    Args:
        svg_path: Path to SVG file
        output_path: Path to save PPTX
        slide_width_inches: PowerPoint slide width in inches (default 10)
        slide_height_inches: PowerPoint slide height in inches (default 7.5)
    """
    svg_path = Path(svg_path)
    output_path = Path(output_path)

    if not svg_path.exists():
        print(f"Error: SVG file not found: {svg_path}")
        sys.exit(1)

    # Parse SVG
    tree = ET.parse(svg_path)
    root = tree.getroot()

    # Get SVG dimensions
    svg_width = float(root.get('width', 1400))
    svg_height = float(root.get('height', 700))

    print(f"SVG dimensions: {svg_width} x {svg_height}")
    print(f"PowerPoint slide: {slide_width_inches}\" x {slide_height_inches}\"")

    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Calculate scaling factors: SVG units → PowerPoint inches
    scale_x = slide_width_inches / svg_width
    scale_y = slide_height_inches / svg_height

    print(f"Scale factors: x={scale_x:.6f}, y={scale_y:.6f}")

    # Process SVG elements
    shapes_added = 0

    # SVG namespace
    ns = {'svg': 'http://www.w3.org/2000/svg'}

    # Process rectangles (boxes)
    for rect in root.findall('.//svg:rect', ns):
        try:
            x = float(rect.get('x', 0))
            y = float(rect.get('y', 0))
            width = float(rect.get('width', 0))
            height = float(rect.get('height', 0))

            # Skip tiny elements
            if width < 2 or height < 2:
                continue

            # Convert to PowerPoint coordinates
            left = Inches(x * scale_x)
            top = Inches(y * scale_y)
            w = Inches(width * scale_x)
            h = Inches(height * scale_y)

            # Add rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, w, h
            )

            # Set fill color
            fill = rect.get('fill', 'none')
            if fill and fill != 'none':
                shape.fill.solid()
                if fill.startswith('#'):
                    rgb = hex_to_rgb(fill)
                    shape.fill.fore_color.rgb = RGBColor(*rgb)
            else:
                shape.fill.background()

            # Set border
            stroke = rect.get('stroke', 'none')
            stroke_width = rect.get('stroke-width', '0')

            if stroke and stroke != 'none':
                shape.line.color.rgb = RGBColor(*hex_to_rgb(stroke))
                try:
                    shape.line.width = Pt(float(stroke_width))
                except:
                    shape.line.width = Pt(1)
            else:
                shape.line.fill.background()

            shapes_added += 1

        except Exception as e:
            print(f"Warning: Failed to process rect: {e}")
            continue

    print(f"✓ Added {shapes_added} rectangles")

    # Process text elements
    texts_added = 0
    for text_elem in root.findall('.//svg:text', ns):
        try:
            x = float(text_elem.get('x', 0))
            y = float(text_elem.get('y', 0))

            # Get text content
            text_content = ''.join(text_elem.itertext()).strip()
            if not text_content:
                continue

            # Get font properties from SVG
            svg_font_size = float(text_elem.get('font-size', 12))
            font_weight = text_elem.get('font-weight', 'normal')
            text_anchor = text_elem.get('text-anchor', 'start')

            # Calculate PowerPoint font size: scale SVG font size to PowerPoint
            # SVG is in pixels, PowerPoint uses points (1/72 inch)
            # Scale from SVG pixels to PowerPoint points
            pptx_font_size = max(6, svg_font_size * scale_y * 10)  # Min 6pt

            # Estimate text box dimensions - generous width to avoid truncation
            char_width = pptx_font_size * 0.7  # Points per character (generous)
            text_width_pts = len(text_content) * char_width * 1.5  # Extra width buffer
            text_height_pts = pptx_font_size * 2.0

            # Convert to inches
            text_width_inches = text_width_pts / 72
            text_height_inches = text_height_pts / 72

            # Convert position to PowerPoint coordinates
            left = Inches(x * scale_x)
            top = Inches(y * scale_y - text_height_inches / 2)  # Center vertically

            # Adjust for text-anchor (alignment)
            if text_anchor == 'middle':
                left -= Inches(text_width_inches / 2)
            elif text_anchor == 'end':
                left -= Inches(text_width_inches)

            w = Inches(text_width_inches)
            h = Inches(text_height_inches)

            # Add text box
            textbox = slide.shapes.add_textbox(left, top, w, h)
            text_frame = textbox.text_frame
            text_frame.text = text_content
            text_frame.word_wrap = False

            # Format text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(pptx_font_size)
                paragraph.font.bold = (font_weight in ['bold', '700', '600'])

                fill = text_elem.get('fill', '#000000')
                if fill.startswith('#'):
                    paragraph.font.color.rgb = RGBColor(*hex_to_rgb(fill))

                # Set alignment
                if text_anchor == 'middle':
                    from pptx.enum.text import PP_ALIGN
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text_anchor == 'end':
                    from pptx.enum.text import PP_ALIGN
                    paragraph.alignment = PP_ALIGN.RIGHT

            texts_added += 1

        except Exception as e:
            print(f"Warning: Failed to process text: {e}")
            continue

    print(f"✓ Added {texts_added} text elements")

    # Now add images from images.json
    images_json_path = svg_path.with_suffix('.images.json')
    if images_json_path.exists():
        images_added = add_images_from_json(slide, images_json_path, scale_x, scale_y)
        print(f"✓ Added {images_added} images")

    # Save presentation
    prs.save(str(output_path))
    print(f"✓ Saved: {output_path}")
    print(f"\nTotal: {shapes_added} shapes, {texts_added} texts")


def add_images_from_json(slide, images_json_path, scale_x, scale_y):
    """Add images from images.json to the slide."""

    with open(images_json_path, 'r') as f:
        data = json.load(f)

    images = data.get('images', [])
    images_added = 0

    for img_data in images:
        try:
            src = img_data['src']
            x = img_data['x']
            y = img_data['y']
            width = img_data['width']
            height = img_data['height']

            # Parse file URL
            if src.startswith('file://'):
                img_path = unquote(urlparse(src).path)
            else:
                img_path = src

            img_path = Path(img_path)

            if not img_path.exists():
                print(f"  Warning: Image not found: {img_path}")
                continue

            # Convert coordinates
            left = Inches(x * scale_x)
            top = Inches(y * scale_y)
            w = Inches(width * scale_x)
            h = Inches(height * scale_y)

            # Add image
            slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
            images_added += 1

        except Exception as e:
            print(f"  Warning: Failed to add image: {e}")
            continue

    return images_added


def main():
    parser = argparse.ArgumentParser(description='Convert SVG to PowerPoint with proper coordinate mapping')
    parser.add_argument('svg_file', help='Input SVG file')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file')
    parser.add_argument('--width', type=float, default=13.33, help='Slide width in inches (default: 13.33 for 16:9)')
    parser.add_argument('--height', type=float, default=7.5, help='Slide height in inches (default: 7.5)')

    args = parser.parse_args()

    svg_to_pptx_custom(args.svg_file, args.output, args.width, args.height)


if __name__ == '__main__':
    main()
