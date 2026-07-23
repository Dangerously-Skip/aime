#!/usr/bin/env python3
"""
Convert SVG to editable PowerPoint shapes using svg2pptx, then add images.

This creates native PowerPoint shapes that can be edited directly in PowerPoint.
"""

import sys
import argparse
import json
from pathlib import Path
from svg2pptx import svg_to_pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from urllib.parse import urlparse, unquote


def convert_svg_to_pptx(svg_path: str, output_path: str):
    """
    Convert SVG to PowerPoint with editable shapes, then add images.

    Args:
        svg_path: Path to SVG file
        output_path: Path to save PPTX
    """
    svg_path = Path(svg_path).resolve()
    output_path = Path(output_path).resolve()
    images_json_path = svg_path.with_suffix('.images.json')

    if not svg_path.exists():
        print(f"Error: SVG file not found: {svg_path}")
        sys.exit(1)

    print(f"Converting SVG to PowerPoint shapes: {svg_path.name}")
    print(f"Output: {output_path}")

    try:
        # Use svg2pptx to convert SVG to PowerPoint shapes
        # This creates native editable shapes
        svg_to_pptx(str(svg_path), str(output_path))
        print(f"✓ SVG shapes converted")

        # Now add images if images.json exists
        if images_json_path.exists():
            print(f"✓ Found images data: {images_json_path.name}")
            add_images_to_pptx(output_path, images_json_path)
        else:
            print(f"⚠ No images.json found - skipping image embedding")

        print(f"✓ Conversion complete!")

    except Exception as e:
        print(f"Error during conversion: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

    print(f"✓ PowerPoint saved: {output_path}")
    print(f"\nAll shapes are now editable in PowerPoint!")
    print("You can:")
    print("  - Move, resize, and recolor shapes")
    print("  - Edit text directly")
    print("  - Ungroup and modify individual elements")

    return output_path


def add_images_to_pptx(pptx_path: Path, images_json_path: Path):
    """Add images from images.json to the PowerPoint file."""

    # Load images data
    with open(images_json_path, 'r') as f:
        images_data = json.load(f)

    images = images_data.get('images', [])
    if not images:
        print("  No images to add")
        return

    print(f"  Adding {len(images)} images...")

    # Load presentation
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]

    # Get slide dimensions from images_data
    svg_width = images_data.get('svg_width', 2100)
    svg_height = images_data.get('svg_height', 1225)

    # PowerPoint slide is 10 inches wide x 7.5 inches tall by default
    pptx_width_inches = 10.0
    pptx_height_inches = 7.5

    # Calculate scaling factors
    scale_x = pptx_width_inches / svg_width
    scale_y = pptx_height_inches / svg_height

    images_added = 0
    for img_data in images:
        try:
            # Parse file:// URL to get local path
            src = img_data['src']
            if src.startswith('file://'):
                img_path = Path(unquote(urlparse(src).path))
            else:
                img_path = Path(src)

            if not img_path.exists():
                print(f"  ⚠ Image not found: {img_path}")
                continue

            # Convert SVG coordinates to PowerPoint inches
            left = Inches(img_data['x'] * scale_x)
            top = Inches(img_data['y'] * scale_y)
            width = Inches(img_data['width'] * scale_x)
            height = Inches(img_data['height'] * scale_y)

            # Add image to slide
            slide.shapes.add_picture(str(img_path), left, top, width, height)
            images_added += 1

        except Exception as e:
            print(f"  ⚠ Failed to add image {img_data.get('src', 'unknown')}: {e}")
            continue

    # Save with images
    prs.save(str(pptx_path))
    print(f"  ✓ Added {images_added}/{len(images)} images")


def main():
    parser = argparse.ArgumentParser(description='Convert SVG to editable PowerPoint shapes')
    parser.add_argument('svg_file', help='Input SVG file')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file')

    args = parser.parse_args()

    convert_svg_to_pptx(args.svg_file, args.output)


if __name__ == '__main__':
    main()
